"""
Build docs/Clinical_Lab_Analysis_System.pptx from docs/SLIDES.md.

Single source of truth: edit SLIDES.md, then regenerate the deck with

    uv run --with python-pptx python docs/build_pptx.py

The first "## Slide" section is rendered as a title slide; the rest become
"Title and Content" slides whose bullets come from the markdown list items.
"""

import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

DOCS = Path(__file__).resolve().parent
SLIDES_MD = DOCS / "SLIDES.md"
OUT_PPTX = DOCS / "Clinical_Lab_Analysis_System.pptx"

SLIDE_HEADER = re.compile(r"^## Slide \d+ - (.+)$")


def parse_sections(md_text):
    """Return a list of (title, [raw_body_lines]) in document order."""
    sections = []
    title, body = None, []
    for line in md_text.splitlines():
        m = SLIDE_HEADER.match(line)
        if m:
            if title is not None:
                sections.append((title, body))
            title, body = m.group(1).strip(), []
        elif title is not None:
            body.append(line)
    if title is not None:
        sections.append((title, body))
    return sections


def clean(text):
    """Drop markdown bold/inline-code markers; keep readable text."""
    return text.replace("**", "").replace("`", "").strip()


def body_to_bullets(body_lines):
    """Collapse a markdown list (with wrapped continuation lines) into bullets."""
    bullets = []
    for raw in body_lines:
        if not raw.strip():
            continue
        if raw.lstrip().startswith("- "):
            bullets.append(clean(raw.lstrip()[2:]))
        elif bullets:  # wrapped continuation of the previous bullet
            bullets[-1] += " " + clean(raw)
        else:
            bullets.append(clean(raw))
    return bullets


def add_title_slide(prs, body_lines):
    lines = [clean(l) for l in body_lines if l.strip()]
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = lines[0] if lines else "Title"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "\n".join(lines[1:])
    return slide


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, text in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = text
        para.level = 0
        for run in para.runs:
            run.font.size = Pt(18)
    return slide


def main():
    sections = parse_sections(SLIDES_MD.read_text(encoding="utf-8"))
    if not sections:
        raise SystemExit("No '## Slide' sections found in SLIDES.md")

    prs = Presentation()
    add_title_slide(prs, sections[0][1])
    for title, body in sections[1:]:
        add_content_slide(prs, title, body_to_bullets(body))

    prs.save(OUT_PPTX)
    print(f"Wrote {OUT_PPTX} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
